"""文档解析器 - 支持多种文件类型。"""

import logging
import warnings
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

from ..core.base import BaseProcessor, ProcessResult
from ..core.schema import DocumentAnnotation, FileType, EXT_TO_FILE_TYPE
from ..core.logger import get_logger

# 抑制 pdfminer 的字体警告
logging.getLogger("pdfminer").setLevel(logging.ERROR)
warnings.filterwarnings("ignore", message=".*FontBBox.*")


class ParserBackend(str, Enum):
    """解析器后端类型。"""
    AUTO = "auto"           # 自动选择（优先 Docling）
    DOCLING = "docling"     # Docling 高精度解析
    LEGACY = "legacy"       # 原有解析器（pdfplumber + python-docx 等）


@dataclass
class DocContent:
    """
    解析后的文档内容。

    Attributes:
        doc_id: 文档标识符
        file_type: 文件类型枚举
        file_path: 文件路径
        page_count: 页数
        text: 提取的文本内容
        pages: 页面图像字节列表（用于PDF/图片）
        metadata: 额外元数据
    """

    doc_id: str
    file_type: FileType
    file_path: str
    page_count: int = 0
    text: str = ""
    pages: List[bytes] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)


def create_parser(backend: ParserBackend = ParserBackend.AUTO, config: Optional[Dict[str, Any]] = None) -> 'BaseProcessor':
    """
    创建文档解析器工厂函数。
    
    Args:
        backend: 解析器后端类型
        config: 配置字典
        
    Returns:
        解析器实例
    """
    logger = get_logger()
    config = config or {}
    
    if backend == ParserBackend.AUTO:
        # 自动选择：优先 Docling
        try:
            from .docling_parser import DoclingParser
            parser = DoclingParser(config)
            if parser.is_available():
                logger.debug("自动选择 Docling 解析器")
                return parser
        except ImportError:
            pass
        logger.debug("Docling 不可用，使用传统解析器")
        return DocParser(config)
    
    elif backend == ParserBackend.DOCLING:
        from .docling_parser import DoclingParser
        return DoclingParser(config)
    
    else:  # LEGACY
        return DocParser(config)


class DocParser(BaseProcessor):
    """
    文档解析器 - 支持多种文件类型。

    支持的格式：
    - PDF: pdfplumber 或 PyPDF2
    - DOC/DOCX: python-docx
    - XLS/XLSX: openpyxl
    - PPT/PPTX: python-pptx
    - HTML/HTM: BeautifulSoup
    - TXT/MD: 直接文本读取
    """

    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """初始化文档解析器。"""
        super().__init__(config)
        self.supported_formats = self.config.get(
            "supported_formats",
            ["pdf", "doc", "docx", "xls", "xlsx", "ppt", "pptx", "html", "htm", "txt", "md"]
        )
        self.extract_images = self.config.get("extract_images", True)
        self.logger = get_logger()

    def process(self, input_data: str) -> ProcessResult:
        """
        解析文档文件。

        Args:
            input_data: 文档文件路径

        Returns:
            包含DocContent的ProcessResult
        """
        file_path = Path(input_data)

        if not file_path.exists():
            self.logger.error(f"文件不存在: {input_data}")
            return ProcessResult(
                success=False,
                errors=[f"文件不存在: {input_data}"]
            )

        file_type = self._detect_file_type(file_path)
        self.logger.parser_start("Legacy (pdfplumber/python-docx)")

        try:
            if file_type == FileType.PDF:
                content = self._parse_pdf(file_path)
            elif file_type == FileType.DOC:
                content = self._parse_docx(file_path)
            elif file_type == FileType.EXCEL:
                content = self._parse_excel(file_path)
            elif file_type == FileType.PPT:
                content = self._parse_pptx(file_path)
            elif file_type in (FileType.HTML,):
                content = self._parse_html(file_path)
            elif file_type in (FileType.TXT, FileType.MD):
                content = self._parse_text(file_path)
            else:
                self.logger.error(f"不支持的文件类型: {file_type}")
                return ProcessResult(
                    success=False,
                    errors=[f"不支持的文件类型: {file_type}"]
                )

            # 记录解析结果
            self.logger.elements_detected(
                page_count=content.page_count,
                images=content.metadata.get("image_count", 1 if content.metadata.get("has_image") else 0),
                tables=content.metadata.get("table_count", 1 if content.metadata.get("has_table") else 0),
                formulas=1 if content.metadata.get("has_formula") else 0,
                charts=content.metadata.get("chart_count", 1 if content.metadata.get("has_chart") else 0),
                table_pages=content.metadata.get("table_pages", []),
                image_pages=content.metadata.get("image_pages", [])
            )

            return ProcessResult(success=True, data=content)

        except Exception as e:
            self.logger.error(f"解析文件时出错 {file_path}: {str(e)}")
            return ProcessResult(
                success=False,
                errors=[f"解析文件时出错 {file_path}: {str(e)}"]
            )

    def _detect_file_type(self, file_path: Path) -> FileType:
        """从扩展名检测文件类型。"""
        ext = file_path.suffix.lower()
        return EXT_TO_FILE_TYPE.get(ext, FileType.TXT)

    def _parse_pdf(self, file_path: Path) -> DocContent:
        """解析PDF文档，提取结构化元素信息。"""
        import io
        import sys
        from contextlib import redirect_stderr

        doc_id = file_path.stem

        # 抑制 pdfminer 的字体警告（直接打印到 stderr）
        _stderr = sys.stderr
        
        # 优先使用pdfplumber（更适合表格和图片）
        try:
            import pdfplumber
            sys.stderr = io.StringIO()

            pages: List[bytes] = []
            text_parts: List[str] = []

            # === 结构化元素检测 ===
            total_images = 0
            total_tables = 0
            total_rects = 0
            total_lines = 0
            total_curves = 0  # 曲线数量（用于判断图表）
            table_pages = set()  # 记录有表格的页码
            image_pages = set()  # 记录有图片的页码
            possible_scanned_table_pages = set()  # 可能是扫描版表格的页码
            
            # 详细表格信息（用于跨页检测）
            tables_detail = []

            with pdfplumber.open(file_path) as pdf:
                for page_idx, page in enumerate(pdf.pages):
                    # 提取文本
                    page_text = page.extract_text() or ""
                    text_parts.append(page_text)

                    # 检测表格
                    page_tables = page.find_tables()
                    table_bboxes = []
                    if page_tables:
                        total_tables += len(page_tables)
                        table_pages.add(page_idx)
                        for tbl_idx, tbl in enumerate(page_tables):
                            bbox = tbl.bbox if hasattr(tbl, 'bbox') else None
                            table_bboxes.append(bbox)
                            # 记录表格详细信息（包含列数用于复杂表格判断）
                            rows = 0
                            cols = 0
                            if hasattr(tbl, 'cells') and tbl.cells:
                                rows = len(set(c[1] for c in tbl.cells if c))  # 根据 top 坐标估算行数
                                cols = len(set(c[0] for c in tbl.cells if c))  # 根据 left 坐标估算列数
                            tables_detail.append({
                                "page": page_idx,
                                "bbox": bbox,
                                "rows": rows,
                                "cols": cols,
                            })
                            self.logger.table_info(tbl_idx, page_idx, rows, cols, bbox=bbox)

                    # 检测图片（排除与表格重叠的区域）
                    page_images = page.images if page.images else []
                    real_images = self._filter_table_images(page_images, table_bboxes)
                    if real_images:
                        total_images += len(real_images)
                        image_pages.add(page_idx)
                        
                        # === 检测可能是扫描版表格的大图片 ===
                        # 如果没有检测到结构化表格，但有占据大部分页面的图片，可能是扫描版
                        if not table_bboxes:  # 该页没有结构化表格
                            page_width = page.width
                            page_height = page.height
                            page_area = page_width * page_height
                            
                            for img in real_images:
                                img_width = img.get('width', 0) or (img.get('x1', 0) - img.get('x0', 0))
                                img_height = img.get('height', 0) or (img.get('y1', 0) - img.get('y0', 0))
                                img_area = img_width * img_height
                                
                                # 如果图片占页面面积 > 50%，可能是扫描版表格
                                if page_area > 0 and img_area / page_area > 0.5:
                                    possible_scanned_table_pages.add(page_idx)
                                    self.logger.debug(
                                        f"页{page_idx}: 发现大图片 ({img_width:.0f}x{img_height:.0f}), "
                                        f"占页面 {img_area/page_area:.1%}，可能是扫描版表格"
                                    )

                    # 检测线条、矩形、曲线
                    page_lines = len(page.lines) if page.lines else 0
                    page_rects = len(page.rects) if page.rects else 0
                    page_curves = len(page.curves) if hasattr(page, 'curves') and page.curves else 0
                    
                    total_lines += page_lines
                    total_rects += page_rects
                    total_curves += page_curves

                    # 如果启用，提取页面图像（用于 OCR 备用）
                    if self.extract_images:
                        img = page.to_image().original
                        img_bytes = io.BytesIO()
                        img.save(img_bytes, format='PNG')
                        pages.append(img_bytes.getvalue())

            # === 改进的 has_chart 判断逻辑 ===
            # 1. 如果有曲线，很可能是图表（折线图、饼图等）
            # 2. 如果矩形/线条很多但没有对应表格，可能是流程图/柱状图
            # 3. 排除表格区域的线条
            has_chart = self._detect_chart(
                total_lines=total_lines,
                total_rects=total_rects,
                total_curves=total_curves,
                total_tables=total_tables,
                page_count=len(pdf.pages)
            )
            
            # 恢复 stderr
            sys.stderr = _stderr
            
            # === 处理可能的扫描版表格（图片表格）===
            # 如果有大图片但该页没有结构化表格，可能是扫描版/截图表格
            has_image_table = len(possible_scanned_table_pages) > 0
            if has_image_table and total_tables == 0:
                self.logger.warning(
                    f"未检测到结构化表格，但页面 {sorted(possible_scanned_table_pages)} 包含大图片，"
                    "可能是扫描版表格。建议使用 Docling 解析器或 OCR 进行识别。"
                )
            
            # === 检测复杂表格 ===
            # 复杂表格判断标准：列数>10 或 表格非常宽（可能有多级表头）
            has_complex_table = self._detect_complex_table(tables_detail)

            return DocContent(
                doc_id=doc_id,
                file_type=FileType.PDF,
                file_path=str(file_path),
                page_count=len(pdf.pages) if not pages else len(pages),
                text="\n".join(text_parts),
                pages=pages,
                metadata={
                    "has_image": total_images > 0,
                    "has_table": total_tables > 0,
                    "has_image_table": has_image_table,
                    "has_complex_table": has_complex_table,
                    "has_formula": False,  # pdfplumber 不检测公式，后续可用 OCR 补充
                    "has_chart": has_chart,
                    "image_count": total_images,
                    "table_count": total_tables,
                    "rect_count": total_rects,
                    "line_count": total_lines,
                    "curve_count": total_curves,
                    "table_pages": sorted(table_pages),
                    "image_pages": sorted(image_pages),
                    "tables_detail": tables_detail,
                    # 扫描版表格提示（向后兼容）
                    "possible_scanned_table": has_image_table and total_tables == 0,
                    "possible_scanned_table_pages": sorted(possible_scanned_table_pages),
                }
            )

        except ImportError:
            # pdfplumber 未安装，恢复 stderr 并回退到 PyPDF2
            sys.stderr = _stderr
            self.logger.parser_fallback("pdfplumber", "PyPDF2", "pdfplumber 未安装")
            try:
                import PyPDF2

                text_parts: List[str] = []

                with open(file_path, "rb") as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    for page in pdf_reader.pages:
                        text_parts.append(page.extract_text() or "")

                return DocContent(
                    doc_id=doc_id,
                    file_type=FileType.PDF,
                    file_path=str(file_path),
                    page_count=len(pdf_reader.pages),
                    text="\n".join(text_parts),
                    pages=[],
                    metadata={
                        "has_image": False,
                        "has_table": False,
                        "has_formula": False,
                        "has_chart": False,
                    }
                )

            except ImportError:
                raise ImportError(
                    "PDF解析需要pdfplumber或PyPDF2。"
                    "请使用: pip install pdfplumber"
                )
    
    def _filter_table_images(self, images: List[Dict], table_bboxes: List) -> List[Dict]:
        """
        过滤掉与表格重叠的图片。
        
        问题：pdfplumber 有时会把表格边框识别为图片。
        解决：如果图片的 bbox 与表格 bbox 高度重叠，则认为是表格而非图片。
        
        Args:
            images: 检测到的图片列表
            table_bboxes: 表格边界框列表
            
        Returns:
            过滤后的真正图片列表
        """
        if not table_bboxes:
            return images
        
        real_images = []
        for img in images:
            img_bbox = (img.get('x0', 0), img.get('top', 0), 
                       img.get('x1', 0), img.get('bottom', 0))
            
            is_table_overlap = False
            for tbl_bbox in table_bboxes:
                if tbl_bbox and self._bbox_overlap(img_bbox, tbl_bbox) > 0.5:
                    is_table_overlap = True
                    self.logger.debug(f"图片与表格重叠，过滤: img={img_bbox}, table={tbl_bbox}")
                    break
            
            if not is_table_overlap:
                real_images.append(img)
        
        return real_images
    
    def _bbox_overlap(self, bbox1: tuple, bbox2: tuple) -> float:
        """
        计算两个边界框的重叠比例。
        
        Args:
            bbox1: (x0, y0, x1, y1)
            bbox2: (x0, y0, x1, y1)
            
        Returns:
            重叠面积占 bbox1 面积的比例
        """
        x0 = max(bbox1[0], bbox2[0])
        y0 = max(bbox1[1], bbox2[1])
        x1 = min(bbox1[2], bbox2[2])
        y1 = min(bbox1[3], bbox2[3])
        
        if x0 >= x1 or y0 >= y1:
            return 0.0
        
        overlap_area = (x1 - x0) * (y1 - y0)
        bbox1_area = (bbox1[2] - bbox1[0]) * (bbox1[3] - bbox1[1])
        
        if bbox1_area <= 0:
            return 0.0
        
        return overlap_area / bbox1_area
    
    def _detect_chart(
        self, 
        total_lines: int, 
        total_rects: int, 
        total_curves: int,
        total_tables: int,
        page_count: int
    ) -> bool:
        """
        改进的图表检测逻辑。
        
        判断依据：
        1. 有曲线（折线图、饼图、圆弧等）→ 很可能是图表
        2. 矩形数量远超表格需要的 → 可能是柱状图
        3. 排除纯表格文档的误判
        
        Args:
            total_lines: 线条总数
            total_rects: 矩形总数
            total_curves: 曲线总数
            total_tables: 表格总数
            page_count: 页数
            
        Returns:
            是否有图表
        """
        # 1. 有曲线 → 很可能是图表（折线图、饼图）
        if total_curves > 5:
            self.logger.debug(f"检测到 {total_curves} 条曲线，判断为有图表")
            return True
        
        # 2. 估算表格需要的线条/矩形数量
        # 一个简单表格大约需要 (rows+1) + (cols+1) 条线
        # 保守估计每个表格平均 20 条线/10 个矩形
        estimated_table_lines = total_tables * 20
        estimated_table_rects = total_tables * 10
        
        # 如果实际线条/矩形远超表格需要，可能有其他图形元素
        extra_lines = total_lines - estimated_table_lines
        extra_rects = total_rects - estimated_table_rects
        
        # 3. 每页平均超过 5 个额外矩形，可能是图表
        if page_count > 0 and extra_rects / page_count > 5:
            self.logger.debug(f"每页平均额外矩形 {extra_rects/page_count:.1f} > 5，可能有图表")
            return True
        
        # 4. 有较多矩形但没有表格，可能是流程图/柱状图
        if total_tables == 0 and total_rects > 10:
            self.logger.debug(f"无表格但有 {total_rects} 个矩形，可能有图表")
            return True
        
        return False
    
    def _detect_complex_table(self, tables_detail: List[Dict]) -> bool:
        """
        检测是否有复杂表格。
        
        复杂表格判断标准（RAG 场景下难以正确解析）：
        1. 列数过多（>10列）- 在 chunk 中难以保持格式
        2. 行数非常多（>100行）- 容易超出 chunk 大小
        3. 表格占据大部分页面 - 可能有多级表头
        
        Args:
            tables_detail: 表格详细信息列表
            
        Returns:
            是否有复杂表格
        """
        if not tables_detail:
            return False
        
        for tbl in tables_detail:
            cols = tbl.get("cols", 0)
            rows = tbl.get("rows", 0)
            
            # 列数过多：> 10 列
            if cols > 10:
                self.logger.debug(f"检测到复杂表格：列数 {cols} > 10")
                return True
            
            # 行数非常多：> 100 行
            if rows > 100:
                self.logger.debug(f"检测到复杂表格：行数 {rows} > 100")
                return True
            
            # 宽表格：列数 >= 7 且行数 > 20
            if cols >= 7 and rows > 20:
                self.logger.debug(f"检测到复杂表格：列数 {cols} >= 7 且行数 {rows} > 20")
                return True
        
        return False

    def _parse_docx(self, file_path: Path) -> DocContent:
        """解析Word文档（.doc, .docx），提取结构化元素信息。"""
        try:
            from docx import Document as DocxDocument
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
            from lxml import etree
        except ImportError:
            raise ImportError(
                "Word文档解析需要python-docx。"
                "请使用: pip install python-docx lxml"
            )

        doc = DocxDocument(file_path)

        # 从段落提取文本
        text_parts = [p.text for p in doc.paragraphs]

        # === 改进的页数估算 ===
        # 方法1: 检测分页符
        page_breaks = 0
        w_ns = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
        for para in doc.paragraphs:
            # 检测 <w:br w:type="page"/>
            breaks = para._element.findall(f'.//{w_ns}br[@{w_ns}type="page"]')
            page_breaks += len(breaks)
            # 检测 <w:pageBreakBefore/>
            if para._element.find(f'.//{w_ns}pageBreakBefore') is not None:
                page_breaks += 1
        
        # === 结构化元素检测（先收集表格信息用于页数估算） ===
        # 1. 检测表格（包含详细信息）
        tables_detail = []
        table_pages = set()
        current_page = 0
        total_table_rows = 0
        
        # 遍历文档的 body 子元素，跟踪页码
        body = doc._element.body
        for child in body:
            # 检测分页符
            if child.tag.endswith('p'):  # 段落
                breaks = child.findall(f'.//{w_ns}br[@{w_ns}type="page"]')
                if breaks:
                    current_page += 1
                if child.find(f'.//{w_ns}pageBreakBefore') is not None:
                    current_page += 1
            
            # 检测表格
            if child.tag.endswith('tbl'):  # 表格
                table_pages.add(current_page)
                # 获取表格行数和列数
                rows = len(child.findall(f'.//{w_ns}tr'))
                # 获取列数（从第一行的单元格数量估算）
                first_row = child.find(f'.//{w_ns}tr')
                cols = len(first_row.findall(f'.//{w_ns}tc')) if first_row is not None else 0
                total_table_rows += rows
                tables_detail.append({
                    "page": current_page,
                    "rows": rows,
                    "cols": cols,
                })
                self.logger.table_info(len(tables_detail)-1, current_page, rows, cols)
        
        # 方法2: 基于内容估算（段落 + 表格行数，每50行约1页）
        total_paragraph_lines = sum(1 for p in doc.paragraphs if p.text.strip())
        total_content_lines = total_paragraph_lines + total_table_rows
        estimated_pages = max(1, total_content_lines // 50 + 1)
        
        # 综合估算
        page_count = max(page_breaks + 1, estimated_pages)
        self.logger.debug(f"DOCX 页数估算: 分页符={page_breaks}, 段落行={total_paragraph_lines}, 表格行={total_table_rows}, 估算页数={page_count}")
        
        has_table = len(doc.tables) > 0

        # 2. 检测图片（内联图片 + 关系中的图片）
        has_image = False
        image_count = 0
        # 检测 inline shapes（内联图片）
        for shape in doc.inline_shapes:
            if shape.type == 3:  # WD_INLINE_SHAPE_TYPE.PICTURE = 3
                has_image = True
                image_count += 1
        # 检测关系中的图片
        if not has_image:
            for rel in doc.part.rels.values():
                if 'image' in rel.reltype:
                    has_image = True
                    image_count += 1

        # 3. 检测公式（OMML - Office Math Markup Language）
        has_formula = False
        omml_ns = '{http://schemas.openxmlformats.org/officeDocument/2006/math}'
        for para in doc.paragraphs:
            if omml_ns in para._element.xml:
                has_formula = True
                break

        # 4. 检测图表（嵌入的 chart）
        has_chart = False
        chart_count = 0
        for rel in doc.part.rels.values():
            if 'chart' in rel.reltype:
                has_chart = True
                chart_count += 1

        # 检测复杂表格
        has_complex_table = self._detect_complex_table(tables_detail)
        
        # 检测图片表格（有图片但没有结构化表格，图片中可能有表格）
        # 对于 DOCX，如果有大量图片但没有表格，可能是扫描版文档
        has_image_table = has_image and not has_table and image_count > 0

        return DocContent(
            doc_id=file_path.stem,
            file_type=FileType.DOC,
            file_path=str(file_path),
            page_count=page_count,
            text="\n".join(text_parts),
            metadata={
                "has_image": has_image,
                "has_table": has_table,
                "has_image_table": has_image_table,
                "has_complex_table": has_complex_table,
                "has_formula": has_formula,
                "has_chart": has_chart,
                "table_count": len(doc.tables),
                "image_count": image_count,
                "chart_count": chart_count,
                "table_pages": sorted(table_pages),
                "tables_detail": tables_detail,
            }
        )

    def _parse_excel(self, file_path: Path) -> DocContent:
        """解析Excel文档（.xls, .xlsx），提取结构化元素信息。"""
        try:
            import openpyxl
        except ImportError:
            raise ImportError(
                "Excel解析需要openpyxl。"
                "请使用: pip install openpyxl"
            )

        # 不使用 read_only 模式以便访问图片和图表
        wb = openpyxl.load_workbook(file_path, read_only=False)

        text_parts: List[str] = []

        # === 结构化元素检测 ===
        has_image = False
        has_chart = False
        has_formula = False
        image_count = 0
        chart_count = 0
        formula_count = 0

        for sheet in wb.worksheets:
            # 提取文本
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    text_parts.append(row_text)

            # 检测图片
            if hasattr(sheet, '_images') and sheet._images:
                has_image = True
                image_count += len(sheet._images)

            # 检测图表
            if hasattr(sheet, '_charts') and sheet._charts:
                has_chart = True
                chart_count += len(sheet._charts)

            # 检测公式（遍历单元格检查是否有公式）
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.data_type == 'f' or (cell.value and isinstance(cell.value, str) and cell.value.startswith('=')):
                        has_formula = True
                        formula_count += 1
                        if formula_count > 10:  # 检测到足够多就停止
                            break
                if formula_count > 10:
                    break

        wb.close()

        return DocContent(
            doc_id=file_path.stem,
            file_type=FileType.EXCEL,
            file_path=str(file_path),
            page_count=len(wb.worksheets),
            text="\n".join(text_parts),
            metadata={
                "has_image": has_image,
                "has_table": True,  # Excel 本身就是表格
                "has_image_table": False,  # Excel 本身是结构化表格
                "has_complex_table": False,  # Excel 表格可直接解析
                "has_formula": has_formula,
                "has_chart": has_chart,
                "image_count": image_count,
                "chart_count": chart_count,
            }
        )

    def _parse_pptx(self, file_path: Path) -> DocContent:
        """解析PowerPoint文档（.ppt, .pptx），提取结构化元素信息。"""
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            raise ImportError(
                "PowerPoint解析需要python-pptx。"
                "请使用: pip install python-pptx"
            )

        prs = Presentation(file_path)

        text_parts: List[str] = []

        # === 结构化元素检测 ===
        has_image = False
        has_table = False
        has_chart = False
        has_formula = False
        table_count = 0
        image_count = 0
        chart_count = 0
        
        # 记录表格/图片所在的页（slide）
        table_pages = set()
        image_pages = set()
        tables_detail = []

        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                # 提取文本
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text)

                # 检测图片
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_image = True
                    image_count += 1
                    image_pages.add(slide_idx)

                # 检测表格
                if shape.has_table:
                    has_table = True
                    table_count += 1
                    table_pages.add(slide_idx)
                    
                    # 获取表格详细信息
                    tbl = shape.table
                    rows = len(tbl.rows) if hasattr(tbl, 'rows') else 0
                    cols = len(tbl.columns) if hasattr(tbl, 'columns') else 0
                    tables_detail.append({
                        "page": slide_idx,  # slide 作为页
                        "rows": rows,
                        "cols": cols,
                    })
                    self.logger.table_info(table_count-1, slide_idx, rows, cols)

                # 检测图表
                if shape.has_chart:
                    has_chart = True
                    chart_count += 1

                # 检测公式（OLE 对象或特定类型）
                if shape.shape_type == MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT:
                    # 可能是公式对象
                    has_formula = True

        # 检测复杂表格
        has_complex_table = self._detect_complex_table(tables_detail)
        
        # 检测图片表格（PPT 中图片可能包含表格截图）
        has_image_table = has_image and not has_table and image_count > 0

        return DocContent(
            doc_id=file_path.stem,
            file_type=FileType.PPT,
            file_path=str(file_path),
            page_count=len(prs.slides),
            text="\n".join(text_parts),
            metadata={
                "has_image": has_image,
                "has_table": has_table,
                "has_image_table": has_image_table,
                "has_complex_table": has_complex_table,
                "has_formula": has_formula,
                "has_chart": has_chart,
                "table_count": table_count,
                "image_count": image_count,
                "chart_count": chart_count,
                "table_pages": sorted(table_pages),
                "image_pages": sorted(image_pages),
                "tables_detail": tables_detail,
            }
        )

    def _parse_html(self, file_path: Path) -> DocContent:
        """解析HTML文档。"""
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError(
                "HTML解析需要beautifulsoup4。"
                "请使用: pip install beautifulsoup4"
            )

        with open(file_path, "r", encoding="utf-8") as f:
            html = f.read()

        soup = BeautifulSoup(html, "html.parser")

        # 移除script和style元素
        for script in soup(["script", "style"]):
            script.decompose()

        # 获取文本
        text = soup.get_text(separator="\n", strip=True)

        return DocContent(
            doc_id=file_path.stem,
            file_type=FileType.HTML,
            file_path=str(file_path),
            page_count=1,
            text=text,
        )

    def _parse_text(self, file_path: Path) -> DocContent:
        """解析纯文本或Markdown文档。"""
        encodings = ["utf-8", "gbk", "gb2312", "latin-1"]

        text = ""
        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    text = f.read()
                break
            except UnicodeDecodeError:
                continue

        # 估算页数（按行数）
        lines = text.split("\n")
        page_count = max(1, len(lines) // 50)

        return DocContent(
            doc_id=file_path.stem,
            file_type=FileType.MD if file_path.suffix == ".md" else FileType.TXT,
            file_path=str(file_path),
            page_count=page_count,
            text=text,
        )

#!/usr/bin/env python
"""创建测试文档。

生成包含各种元素的测试文档用于验证标注功能。
"""

import sys
sys.path.insert(0, "docs_annotation/src")


def create_docx_with_tables():
    """创建包含表格的 DOCX 文档。"""
    try:
        from docx import Document
        from docx.shared import Inches
    except ImportError:
        print("需要 python-docx: pip install python-docx")
        return None
    
    doc = Document()
    
    # 添加标题
    doc.add_heading('测试文档 - 表格跨页', 0)
    
    # 添加一些段落
    doc.add_paragraph('这是一个用于测试跨页表格检测的文档。')
    doc.add_paragraph('下面是一个大表格，可能会跨页。')
    
    # 创建一个大表格（模拟跨页）
    table = doc.add_table(rows=50, cols=4)
    table.style = 'Table Grid'
    
    # 填充表头
    header_cells = table.rows[0].cells
    header_cells[0].text = '序号'
    header_cells[1].text = '名称'
    header_cells[2].text = '数值'
    header_cells[3].text = '备注'
    
    # 填充数据
    for i in range(1, 50):
        row_cells = table.rows[i].cells
        row_cells[0].text = str(i)
        row_cells[1].text = f'项目{i}'
        row_cells[2].text = str(i * 100)
        row_cells[3].text = f'备注信息{i}'
    
    # 添加更多内容
    doc.add_paragraph('')
    doc.add_paragraph('表格结束。')
    
    # 添加第二个表格
    doc.add_heading('第二部分', level=1)
    table2 = doc.add_table(rows=10, cols=3)
    table2.style = 'Table Grid'
    
    for i in range(10):
        row = table2.rows[i].cells
        row[0].text = f'A{i}'
        row[1].text = f'B{i}'
        row[2].text = f'C{i}'
    
    # 保存
    output_path = 'docs_annotation/test/test_data/test_cross_page_table.docx'
    import os
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    doc.save(output_path)
    print(f"✅ 创建: {output_path}")
    return output_path


def create_xlsx_with_data():
    """创建包含数据的 Excel 文档。"""
    try:
        import openpyxl
    except ImportError:
        print("需要 openpyxl: pip install openpyxl")
        return None
    
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "数据表"
    
    # 添加表头
    ws['A1'] = '序号'
    ws['B1'] = '名称'
    ws['C1'] = '数值'
    ws['D1'] = '状态'
    
    # 添加数据
    for i in range(2, 102):
        ws[f'A{i}'] = i - 1
        ws[f'B{i}'] = f'项目{i-1}'
        ws[f'C{i}'] = (i - 1) * 10
        ws[f'D{i}'] = '正常' if i % 2 == 0 else '异常'
    
    # 保存
    output_path = 'docs_annotation/test/test_data/test_excel.xlsx'
    import os
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    wb.save(output_path)
    print(f"✅ 创建: {output_path}")
    return output_path


def create_pptx_with_tables():
    """创建包含表格的 PPT 文档。"""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("需要 python-pptx: pip install python-pptx")
        return None
    
    prs = Presentation()
    
    # 幻灯片 1: 标题
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "测试 PPT"
    subtitle.text = "包含表格的演示文稿"
    
    # 幻灯片 2: 表格
    slide_layout = prs.slide_layouts[5]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加表格
    rows, cols = 5, 4
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(3)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 填充表头
    table.cell(0, 0).text = '序号'
    table.cell(0, 1).text = '名称'
    table.cell(0, 2).text = '数值'
    table.cell(0, 3).text = '状态'
    
    # 填充数据
    for i in range(1, rows):
        table.cell(i, 0).text = str(i)
        table.cell(i, 1).text = f'项目{i}'
        table.cell(i, 2).text = str(i * 100)
        table.cell(i, 3).text = '完成'
    
    # 幻灯片 3: 另一个表格
    slide = prs.slides.add_slide(slide_layout)
    table2 = slide.shapes.add_table(3, 3, left, top, width, Inches(2)).table
    for i in range(3):
        for j in range(3):
            table2.cell(i, j).text = f'({i},{j})'
    
    # 保存
    output_path = 'docs_annotation/test/test_data/test_ppt.pptx'
    import os
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"✅ 创建: {output_path}")
    return output_path


def main():
    """创建所有测试文档。"""
    print("创建测试文档...")
    print("=" * 50)
    
    files = []
    
    # 创建 DOCX
    path = create_docx_with_tables()
    if path:
        files.append(path)
    
    # 创建 Excel
    path = create_xlsx_with_data()
    if path:
        files.append(path)
    
    # 创建 PPT
    path = create_pptx_with_tables()
    if path:
        files.append(path)
    
    print("=" * 50)
    print(f"共创建 {len(files)} 个测试文件")
    
    return files


if __name__ == "__main__":
    main()
